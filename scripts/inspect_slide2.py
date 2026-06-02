import sys
from pptx import Presentation

sys.stdout.reconfigure(encoding='utf-8')
prs = Presentation("NL_Claude_TrustEvaluators.pptx")

# Slide 2 is index 1
slide_2 = prs.slides[1]
print(f"Slide 2: shapes count: {len(slide_2.shapes)}")
for i, s in enumerate(slide_2.shapes):
    txt = s.text_frame.text.strip().replace('\n', ' ') if s.has_text_frame else "No Text"
    print(f"Shape {i}: Name='{s.name}', Type={s.shape_type}, Text='{txt[:80]}'")
