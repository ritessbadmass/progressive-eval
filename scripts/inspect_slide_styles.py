import sys
from pptx import Presentation

# Reconfigure stdout to use UTF-8
sys.stdout.reconfigure(encoding='utf-8')

prs = Presentation("NL_Claude_TrustEvaluators.pptx")

slide = prs.slides[2] # Slide 3 (Research Findings)
print(f"Slide 3 Background: {slide.background}")
for j, shape in enumerate(slide.shapes):
    if shape.has_text_frame:
        print(f"Shape {j}: Name={shape.name}")
        print(f"  Left={shape.left}, Top={shape.top}, Width={shape.width}, Height={shape.height}")
        if shape.fill.type is not None:
            print(f"  Fill Type={shape.fill.type}, Color={shape.fill.fore_color.rgb if hasattr(shape.fill.fore_color, 'rgb') else 'None'}")
        if shape.line.fill.type is not None:
            print(f"  Line Color={shape.line.color.rgb if hasattr(shape.line.color, 'rgb') else 'None'}, Width={shape.line.width}")
        # Inspect first paragraph font
        if len(shape.text_frame.paragraphs) > 0:
            p = shape.text_frame.paragraphs[0]
            if len(p.runs) > 0:
                run = p.runs[0]
                print(f"  Font: Name={run.font.name}, Size={run.font.size}, Color={run.font.color.rgb if hasattr(run.font.color, 'rgb') else 'None'}")
            else:
                print(f"  Font: Name={p.font.name}, Size={p.font.size}, Color={p.font.color.rgb if hasattr(p.font.color, 'rgb') else 'None'}")
print("Done styling inspection")
