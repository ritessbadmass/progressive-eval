import sys
from pptx import Presentation

# Reconfigure stdout to use UTF-8
sys.stdout.reconfigure(encoding='utf-8')

prs = Presentation("NL_Claude_TrustEvaluators.pptx")

with open("scripts/slide_full_text_updated.txt", "w", encoding="utf-8") as f:
    f.write(f"Total Slides: {len(prs.slides)}\n")
    for i, slide in enumerate(prs.slides):
        f.write(f"\n=========================================\n")
        f.write(f"--- Slide {i+1} ---\n")
        f.write(f"=========================================\n")
        title = slide.shapes.title.text if slide.shapes.title else "No Title"
        f.write(f"TITLE: {title}\n\n")
        for j, shape in enumerate(slide.shapes):
            f.write(f"Shape {j}: Name='{shape.name}', Type={shape.shape_type}\n")
            if shape.has_text_frame:
                f.write(f"  TEXT:\n")
                for p in shape.text_frame.paragraphs:
                    if p.text.strip():
                        f.write(f"    - {p.text.strip()}\n")
            # If shape is a table, print cell texts!
            if shape.has_table:
                f.write(f"  TABLE DETAILS:\n")
                for r_idx, row in enumerate(shape.table.rows):
                    row_texts = [cell.text.strip().replace('\n', ' ') for cell in row.cells]
                    f.write(f"    - Row {r_idx}: {row_texts}\n")
            f.write("\n")
print("Done writing to scripts/slide_full_text_updated.txt")
