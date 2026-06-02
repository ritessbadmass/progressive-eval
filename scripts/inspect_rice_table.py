import sys
from pptx import Presentation

sys.stdout.reconfigure(encoding='utf-8')
prs = Presentation("NL_Claude_TrustEvaluators.pptx")

# Slide 6 is index 5
slide = prs.slides[5]
print(f"Slide index 5 Title: {slide.shapes.title.text if slide.shapes.title else 'No Title'}")
for i, s in enumerate(slide.shapes):
    if s.has_text_frame:
        txt = s.text_frame.text.strip()
        if any(keyword in txt for keyword in ["Task Lens", "Interactive Decision", "Progressive", "RICE", "Reach", "Impact", "Confidence", "Effort"]):
            print(f"Shape {i}: Text='{txt}' | Left={s.left}, Top={s.top}, Width={s.width}, Height={s.height}")
        elif len(txt) <= 5 and txt != "":
            print(f"Shape {i} (Short text): Text='{txt}' | Left={s.left}, Top={s.top}, Width={s.width}, Height={s.height}")
