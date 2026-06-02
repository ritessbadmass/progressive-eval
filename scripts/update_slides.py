import sys
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor

# Reconfigure stdout to use UTF-8
sys.stdout.reconfigure(encoding='utf-8')

print("Loading presentation...")
prs = Presentation("NL_Claude_TrustEvaluators.pptx")

# Define helper functions
def move_slide(prs, old_index, new_index):
    xml_slides = prs.slides._sldIdLst
    slide_id_element = xml_slides[old_index]
    xml_slides.remove(slide_id_element)
    xml_slides.insert(new_index, slide_id_element)

def delete_slide(prs, index):
    id_list = prs.slides._sldIdLst
    del id_list[index]

# ==========================================
# STEP 1: MODIFTY SLIDE 1 (TITLE & LANDSCAPE)
# ==========================================
print("Modifying Slide 1...")
slide_1 = prs.slides[0]

# Shapes to delete on Slide 1 for the old trust gap columns
names_to_delete = [
    "Shape 17", "Text 18", "Text 19", "Text 20", "Text 21", "Text 22", 
    "Text 23", "Text 24", "Text 25", "Text 26", "Text 27", "Text 28"
]

shapes_to_remove = [s for s in slide_1.shapes if s.name in names_to_delete]
print(f"Removing {len(shapes_to_remove)} old trust gap shapes from Slide 1...")
for s in shapes_to_remove:
    slide_1.shapes._spTree.remove(s._element)

# Add compact 4-row landscape table at the bottom of Slide 1
# Coordinates in EMUs (Left: 0.5 in, Top: 4.6 in, Width: 12.33 in, Height: 2.1 in)
left = 457200
top = 4434840
width = 11247120
height = 2000000

print("Adding compact 4-row landscape table to Slide 1...")
table_shape = slide_1.shapes.add_table(5, 4, left, top, width, height)
table = table_shape.table

# Column widths
table.columns[0].width = int(width * 0.23)
table.columns[1].width = int(width * 0.23)
table.columns[2].width = int(width * 0.34)
table.columns[3].width = int(width * 0.20)

# Table texts
table_data = [
    ["Solution Type", "Focus Area", "Core Failure / Limitation", "Perceived vs. Actual Quality Gap"],
    ["Basic Citations & Links\n(Perplexity, Gemini)", "Web source attribution & raw links", "Convincing URLs are easily spoofed; users rarely click through to verify.", "High perceived quality, but hides stale data and source mismatches."],
    ["Fact-Checkers\n(Hallucination tools)", "Factual validity checks", "Treats trust as binary correctness; ignores logic, code safety, or tone.", "Mismatches cognitive load, leading to alert fatigue."],
    ["Opaque \"Trust Scores\"\n(Generic badges/scores)", "Percentage-based indicators", "Encourages authority bias; users stop reading the text and trust the score.", "Promotes cognitive laziness rather than critical thinking."],
    ["Progressive Evaluation\n(Our Selected Concept)", "Contextual peer-agent lenses on demand", "Requires active opt-in; no auto-pass grades.", "Calibrates trust at the exact moment of decision."]
]

for r_idx, row in enumerate(table.rows):
    for c_idx, cell in enumerate(row.cells):
        cell.text = table_data[r_idx][c_idx]
        
        # Color styling
        cell.fill.solid()
        if r_idx == 0:
            cell.fill.fore_color.rgb = RGBColor(30, 58, 95) # 1E3A5F (Header slate)
        elif r_idx == 4:
            cell.fill.fore_color.rgb = RGBColor(18, 38, 62) # Highlighted Row background
        else:
            cell.fill.fore_color.rgb = RGBColor(13, 27, 42) # 0D1B2A (matching slide background)
        
        # Padding
        cell.margin_left = 100000
        cell.margin_right = 100000
        cell.margin_top = 80000
        cell.margin_bottom = 80000
        
        # Format text frame and runs
        from pptx.enum.text import PP_ALIGN
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER if r_idx == 0 else PP_ALIGN.LEFT
            for run in p.runs:
                run.font.name = "Arial"
                run.font.size = Pt(10.5)
                if r_idx == 0:
                    run.font.bold = True
                    run.font.color.rgb = RGBColor(255, 255, 255) # White text
                elif r_idx == 4:
                    # Highlight our winning concept in teal/amber
                    run.font.bold = (c_idx == 0)
                    run.font.color.rgb = RGBColor(0, 180, 216) # 00B4D8 teal
                else:
                    run.font.color.rgb = RGBColor(240, 240, 240) # White/light gray
                    if c_idx == 0:
                        run.font.bold = True

# ==========================================
# STEP 2: DELETE ORIGINAL SLIDE 2 (LANDSCAPE)
# ==========================================
print("Deleting original Slide 2 (Landscape)...")
delete_slide(prs, 1)

# Current slide indices shift:
# 0: S1 (Title & Landscape)
# 1: S3 (Research Findings)
# 2: S4 (Target Segment)
# 3: S5 (User Journey Map)
# 4: S6 (Ideation / RICE)
# 5: S7 (Why Progressive Evaluation wins)
# 6: S8 (Product Flow)
# 7: S9 (6 Specialized Evaluators)
# 8: S10 (Success & Risks)

# ==========================================
# STEP 3: REORDER SLIDE 5 (WHY PROGRESSIVE EVALUATION WINS) TO SLIDE 4
# ==========================================
print("Moving 'Why Progressive Evaluation wins' to Slide 4...")
# S7 is currently index 5
move_slide(prs, 5, 3)

# Current slide indices shift:
# 0: S1 (Title & Landscape) -> Slide 1
# 1: S3 (Research Findings) -> Slide 2
# 2: S4 (Target Segment) -> Slide 3
# 3: S7 (Why Progressive Evaluation wins) -> Slide 4
# 4: S5 (User Journey Map) -> Slide 5 (temp)
# 5: S6 (Ideation / RICE) -> Slide 6 (temp)
# 6: S8 (Product Flow) -> Slide 7
# 7: S9 (6 Specialized Evaluators) -> Slide 8
# 8: S10 (Success & Risks) -> Slide 9

# ==========================================
# STEP 4: ADD NEW SLIDE 5 (PROBLEM FRAMING CANVAS)
# ==========================================
print("Creating Slide 5 (Problem Framing Canvas)...")
# Insert a new slide at the end
new_slide = prs.slides.add_slide(prs.slide_layouts[0])

# Re-index new slide to index 4 (so it becomes Slide 5, shifting S5 to index 5 and S6 to index 6)
move_slide(prs, 9, 4)

# Current slide indices are now:
# 0: S1 (Title & Landscape) -> Slide 1
# 1: S3 (Research Findings) -> Slide 2
# 2: S4 (Target Segment) -> Slide 3
# 3: S7 (Why Progressive Evaluation wins) -> Slide 4
# 4: NewSlide (Problem Framing Canvas) -> Slide 5
# 5: S5 (User Journey Map) -> Slide 6
# 6: S6 (Ideation / RICE) -> Slide 7
# 7: S8 (Product Flow) -> Slide 8
# 8: S9 (6 Specialized Evaluators) -> Slide 9
# 9: S10 (Success & Risks) -> Slide 10

# Style new Slide 5 background
bg_rect = new_slide.shapes.add_shape(1, 0, 0, 12161520, 6858000) # Rectangle
bg_rect.fill.solid()
bg_rect.fill.fore_color.rgb = RGBColor(13, 27, 42) # 0D1B2A (Navy)
bg_rect.line.fill.background()

# Style new Slide 5 top bar
top_bar = new_slide.shapes.add_shape(1, 0, 0, 12161520, 804672)
top_bar.fill.solid()
top_bar.fill.fore_color.rgb = RGBColor(13, 27, 42) # 0D1B2A
top_bar.line.fill.background()

# Title text frame (Synthesizing Statement)
title_box = new_slide.shapes.add_textbox(457200, 73152, 11247120, 658368)
tf = title_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "To bridge the judgment gap, Claude must transition from a silent execution engine to a transparent partner in cognitive calibration."
p.font.name = "Georgia"
p.font.size = Pt(17)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)

# Rearrange shapes in shape list so background sits at index 0 (Behind other shapes)
new_slide.shapes._spTree.remove(bg_rect._element)
new_slide.shapes._spTree.insert(2, bg_rect._element) # 0: top_bar, 1: title_box, 2: bg_rect

# Cards Data for 2x3 grid (2 top row, 3 bottom row)
cards_data = [
    {
        "header": "WHAT IS THE TRUE PROBLEM?",
        "body": "No structured evaluation layer exists between a user's intent and Claude's execution. Claude decides silently which capability to use. The user never sees it, never learns from it, and loses agency. Claude was built for speed, and silently removed user agency in the process.",
        "left": 457200, "top": 1200000, "width": 5394960, "height": 2200000
    },
    {
        "header": "WHO FACES IT?",
        "body": "Working professionals using Claude for high-stakes tasks daily — who need to move fast but cannot afford subtle quality mistakes.\n\n→ 71% of surveyed users rate their tasks as high-stakes\n→ 85% use AI tools daily or multiple times a day",
        "left": 6309360, "top": 1200000, "width": 5394960, "height": 2200000
    },
    {
        "header": "HOW DO WE KNOW IT'S REAL?",
        "body": "30% over-trust (Fluency Bias trap)\n37% skeptic trap (rechecking endlessly)\nErrors are discovered only at the point of action.\n\n*Grounded in: 35 surveys + 6 interviews + 2 task observation sessions*",
        "left": 457200, "top": 3750000, "width": 3482373, "height": 2550000
    },
    {
        "header": "WHAT VALUE IS GENERATED?",
        "body": "For Users: Calibrated confidence in outputs and significant speedups in manual auditing.\n\nFor Anthropic: Long-term retention, high-stakes workflow stickiness, and higher Pro/Team upgrades.",
        "left": 4339573, "top": 3750000, "width": 3482373, "height": 2550000
    },
    {
        "header": "WHY SHOULD WE SOLVE NOW?",
        "body": "Trust is a habit. Every session a user spends defaulting to blind copy-pasting or basic prompts reinforces bad habits. The longer we delay introducing structured evaluation, the harder it gets to change user behaviors.",
        "left": 8221946, "top": 3750000, "width": 3482373, "height": 2550000
    }
]

print("Adding 5 rounded rectangle cards to Slide 5...")
for card in cards_data:
    box = new_slide.shapes.add_shape(5, card["left"], card["top"], card["width"], card["height"]) # rounded rect
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(30, 58, 95) # 1E3A5F (Dark slate blue)
    box.line.color.rgb = RGBColor(51, 65, 85) # 334155 border
    box.line.width = Pt(1)
    
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 180000
    tf.margin_right = 180000
    tf.margin_top = 150000
    tf.margin_bottom = 150000
    
    p0 = tf.paragraphs[0]
    p0.text = card["header"]
    p0.font.name = "Arial"
    p0.font.size = Pt(11)
    p0.font.bold = True
    p0.font.color.rgb = RGBColor(0, 180, 216) # 00B4D8 teal
    
    p1 = tf.add_paragraph()
    p1.text = card["body"]
    p1.font.name = "Arial"
    p1.font.size = Pt(10)
    p1.font.color.rgb = RGBColor(240, 240, 240)
    p1.space_before = Pt(4)

# ==========================================
# STEP 5: UPDATE SLIDE 7 (IDEATION / RICE TABLE)
# ==========================================
print("Modifying Slide 7 (Ideation & RICE updates)...")
slide_7 = prs.slides[6]

# 1. Update Idea Wording (What it is / Promise / Limit)
# Idea 1 (Task Lens) Text Boxes:
# Shape 14: WHAT IT IS body
# Shape 16: PROMISE body
# Shape 18: LIMIT body
slide_7.shapes[14].text_frame.text = "Claude infers the task type from the user's prompt and tailors both the answer and the review criteria to it (e.g. coding, research, or career)."
slide_7.shapes[16].text_frame.text = "Zero upfront friction. Keeps Claude's speed-preserving core while offering a post-answer review chip."
slide_7.shapes[18].text_frame.text = "Still depends on Claude correctly inferring the task lens, and the review format can feel slightly abstract."

# Idea 2 (Decision Paths) Text Boxes:
# Shape 25: WHAT IT IS body
# Shape 27: PROMISE body
# Shape 29: LIMIT body
slide_7.shapes[25].text_frame.text = "Surfaced only on ambiguous/high-stakes prompts. Claude previews 2–3 possible ways it could approach the task (e.g. fast, risk-aware)."
slide_7.shapes[27].text_frame.text = "Allows users to steer the response strategy before Claude commits, making the reasoning strategy visible."
slide_7.shapes[29].text_frame.text = "Still introduces a moment of hesitation before users see value. Risks feeling fragile in conversational speed."

# Idea 3 (Progressive Evaluation) Text Boxes:
# Shape 36: WHAT IT IS body
# Shape 38: PROMISE body
# Shape 40: LIMIT body
slide_7.shapes[36].text_frame.text = "Claude answers first at full speed, then offers an optional evaluation layer where users route the answer through 6 specialized peer-evaluators (Code, Research, Writing, Reasoning, Career, Risk) and 3 playbooks (Balanced, Rigor, Style) powered by Llama 3.1 70B NIM models."
slide_7.shapes[38].text_frame.text = "Opt-in = no speed friction by default. Outputs transparent checklists showing strengths, weaknesses, and a concrete 'Verify before use' manual checklist. No opaque trust scores."
slide_7.shapes[40].text_frame.text = "Slightly higher engineering effort than simple review chips because it requires evaluator orchestration, parallel NIM routing, and UI safeguards."

# Format runs on text boxes
text_box_indices = [14, 16, 18, 25, 27, 29, 36, 38, 40]
for idx in text_box_indices:
    tf = slide_7.shapes[idx].text_frame
    tf.word_wrap = True
    for p in tf.paragraphs:
        for run in p.runs:
            run.font.name = "Arial"
            run.font.size = Pt(10)
            run.font.color.rgb = RGBColor(240, 240, 240)

# 2. Update RICE Table Cell Values
rice_updates = {
    # Task Lens + Review Mode (Shape 50)
    51: "7.0",   # Reach
    52: "2.5",   # Impact
    53: "75%",   # Confidence
    54: "4.0",   # Effort
    55: "3.28",  # RICE Score
    
    # Interactive Decision Paths (Shape 57)
    58: "6.5",   # Reach
    59: "2.25",  # Impact
    60: "70%",   # Confidence
    61: "3.5",   # Effort
    62: "2.93",  # RICE Score
    
    # Progressive Eval + Evaluators (Shape 64)
    65: "8.5",   # Reach
    66: "3.0",   # Impact
    67: "85%",   # Confidence
    68: "4.5",   # Effort
    69: "4.82"   # RICE Score
}

print("Updating RICE Table cell shapes...")
from pptx.enum.text import PP_ALIGN
for idx, new_val in rice_updates.items():
    shape = slide_7.shapes[idx]
    shape.text_frame.text = new_val
    for p in shape.text_frame.paragraphs:
        p.alignment = PP_ALIGN.CENTER # Center align
        for run in p.runs:
            run.font.name = "Arial"
            run.font.size = Pt(11)
            # Make the winner scores highlighted or bold
            if idx in [65, 66, 67, 68, 69]:
                run.font.bold = True
                run.font.color.rgb = RGBColor(0, 180, 216) # Teal winner!
            else:
                run.font.color.rgb = RGBColor(255, 255, 255) # White

# Save updated PowerPoint file
print("Saving modified presentation...")
prs.save("NL_Claude_TrustEvaluators.pptx")
print("PowerPoint file successfully updated and saved!")
