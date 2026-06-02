import sys
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Reconfigure stdout to use UTF-8
sys.stdout.reconfigure(encoding='utf-8')

print("Loading presentation...")
prs = Presentation("NL_Claude_TrustEvaluators.pptx")

# ==========================================
# STEP 1: MODIFY SLIDE 1 (TITLE & COMPACT TABLE)
# ==========================================
print("Modifying Slide 1...")
slide_1 = prs.slides[0]

# Shapes to delete on Slide 1 for the old trust gap columns (shapes 17 to 28) and any old tables
names_to_delete = [
    "Shape 17", "Text 18", "Text 19", "Text 20", "Text 21", "Text 22", 
    "Text 23", "Text 24", "Text 25", "Text 26", "Text 27", "Text 28"
]

shapes_to_remove_s1 = [
    s for s in slide_1.shapes 
    if s.name in names_to_delete or s.has_table or s.name.startswith("Table")
]
print(f"Removing {len(shapes_to_remove_s1)} old trust gap / table shapes from Slide 1...")
for s in shapes_to_remove_s1:
    slide_1.shapes._spTree.remove(s._element)

# Add compact 4-row landscape table to the bottom of Slide 1
left = 457200
top = 4434840
width = 11247120
height = 2000000

print("Adding compact 4-row landscape table to Slide 1...")
table_shape = slide_1.shapes.add_table(5, 4, left, top, width, height)
table = table_shape.table

# Column widths
table.columns[0].width = int(width * 0.23)
table.columns[1].width = int(width * 0.23)
table.columns[2].width = int(width * 0.34)
table.columns[3].width = int(width * 0.20)

table_data = [
    ["Solution Type", "Focus Area", "Core Failure / Limitation", "Perceived vs. Actual Quality Gap"],
    ["Basic Citations & Links\n(Perplexity, Gemini)", "Web source attribution & raw links", "Convincing URLs are easily spoofed; users rarely click through to verify.", "High perceived quality, but hides stale data and source mismatches."],
    ["Fact-Checkers\n(Hallucination tools)", "Factual validity checks", "Treats trust as binary correctness; ignores logic, code safety, or tone.", "Mismatches cognitive load, leading to alert fatigue."],
    ["Opaque \"Trust Scores\"\n(Generic badges/scores)", "Percentage-based indicators", "Encourages authority bias; users stop reading the text and trust the score.", "Promotes cognitive laziness rather than critical thinking."],
    ["Progressive Evaluation\n(Our Selected Concept)", "Contextual peer-agent lenses on demand", "Requires active opt-in; no auto-pass grades.", "Calibrates trust at the exact moment of decision."]
]

for r_idx, row in enumerate(table.rows):
    for c_idx, cell in enumerate(row.cells):
        cell.text = table_data[r_idx][c_idx]
        cell.fill.solid()
        
        # Row coloring
        if r_idx == 0:
            cell.fill.fore_color.rgb = RGBColor(30, 58, 95) # 1E3A5F Slate blue
        elif r_idx == 4:
            cell.fill.fore_color.rgb = RGBColor(18, 38, 62) # Highlight Row background
        else:
            cell.fill.fore_color.rgb = RGBColor(13, 27, 42) # 0D1B2A Slide background
        
        # Margins/Padding
        cell.margin_left = 100000
        cell.margin_right = 100000
        cell.margin_top = 80000
        cell.margin_bottom = 80000
        
        # Format text frame and runs
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER if r_idx == 0 else PP_ALIGN.LEFT
            for run in p.runs:
                run.font.name = "Arial"
                run.font.size = Pt(10.5)
                if r_idx == 0:
                    run.font.bold = True
                    run.font.color.rgb = RGBColor(255, 255, 255)
                elif r_idx == 4:
                    run.font.bold = (c_idx == 0)
                    run.font.color.rgb = RGBColor(0, 180, 216) # glowing teal 00B4D8
                else:
                    run.font.color.rgb = RGBColor(240, 240, 240)
                    if c_idx == 0:
                        run.font.bold = True

# ==========================================
# STEP 2: REPURPOSE SLIDE 2 (MARKET LANDSCAPE -> PROBLEM FRAMING CANVAS FIDELITY)
# ==========================================
print("Repurposing Slide 2 (index 1) to become the Problem Framing Canvas with correct content...")
slide_2 = prs.slides[1]

# Clear all existing shapes completely on Slide 2
print(f"Clearing all {len(slide_2.shapes)} existing shapes on Slide 2...")
while len(slide_2.shapes) > 0:
    slide_2.shapes._spTree.remove(slide_2.shapes[0]._element)

# 1. Add Slide 2 Background: Light Warm Cream (Hex: F4F0E6, RGB: 244, 240, 230)
bg_rect = slide_2.shapes.add_shape(1, 0, 0, 12161520, 6858000) # Rectangle
bg_rect.fill.solid()
bg_rect.fill.fore_color.rgb = RGBColor(244, 240, 230)
bg_rect.line.fill.background() # No border

# 2. Add Title Box: Warm Asterisk and Editorial title
title_box = slide_2.shapes.add_textbox(457200, 300000, 11247120, 600000)
tf = title_box.text_frame
tf.word_wrap = True
tf.margin_left = 0
tf.margin_right = 0
tf.margin_top = 0
tf.margin_bottom = 0
p = tf.paragraphs[0]

# Orange pinwheel asterisk ✳
r0 = p.add_run()
r0.text = "✳  "
r0.font.name = "Georgia"
r0.font.size = Pt(23)
r0.font.bold = True
r0.font.color.rgb = RGBColor(195, 106, 79) # Terracotta #C36A4F

# Title Text
r1 = p.add_run()
r1.text = "No Layer Between Intent and Output."
r1.font.name = "Georgia"
r1.font.size = Pt(23)
r1.font.bold = True
r1.font.color.rgb = RGBColor(25, 25, 25) # Dark Charcoal #191919

# Helper Function to add structured bold text blocks
def add_formatted_text(paragraph, segments, font_size=10.5):
    for text, is_bold in segments:
        run = paragraph.add_run()
        run.text = text
        run.font.name = "Arial"
        run.font.size = Pt(font_size)
        run.font.color.rgb = RGBColor(25, 25, 25) # Dark charcoal
        run.font.bold = is_bold

# Grid Coordinates
top_row_y = 1100000
top_row_h = 2400000
card_w = 3549040
col_gap = 300000

bottom_row_y = 3750000
bottom_row_h = 2650000

# CARD 1: What is the True Problem? (Top Left)
c1_left = 457200
# Terracotta box background
box_1 = slide_2.shapes.add_shape(5, c1_left, top_row_y, card_w, top_row_h) # Rounded rect
box_1.fill.solid()
box_1.fill.fore_color.rgb = RGBColor(195, 106, 79) # Terracotta
box_1.line.fill.background()
# Card Title Box (Terrocotta title text)
t_box_1 = slide_2.shapes.add_textbox(c1_left + 150000, top_row_y + 80000, card_w - 300000, 350000)
t_box_1.text_frame.word_wrap = True
t_box_1.text_frame.margin_left = 0
t_box_1.text_frame.margin_top = 0
p_t1 = t_box_1.text_frame.paragraphs[0]
p_t1.text = "What is the True Problem?"
p_t1.font.name = "Arial"
p_t1.font.size = Pt(12)
p_t1.font.bold = True
p_t1.font.color.rgb = RGBColor(25, 25, 25)
# Inner White Card
w_card_1 = slide_2.shapes.add_shape(5, c1_left + 150000, top_row_y + 480000, card_w - 300000, top_row_h - 630000)
w_card_1.fill.solid()
w_card_1.fill.fore_color.rgb = RGBColor(255, 255, 255) # Pure White
w_card_1.line.color.rgb = RGBColor(195, 106, 79) # Terracotta border
w_card_1.line.width = Pt(1)
tf_1 = w_card_1.text_frame
tf_1.word_wrap = True
tf_1.margin_left = 120000
tf_1.margin_right = 120000
tf_1.margin_top = 100000
tf_1.margin_bottom = 100000
p_b1 = tf_1.paragraphs[0]
segments_1 = [
    ("No structured evaluation layer", True),
    (" exists between a user's intent and Claude's execution. Claude ", False),
    ("decides silently", True),
    (" which capability or reasoning path to use. The user never sees it, never learns from it, and loses agency. Claude was built for speed, and silently ", False),
    ("removed user agency", True),
    (" in the process.", False)
]
add_formatted_text(p_b1, segments_1)
# Small Link Icon at bottom right
link_1 = slide_2.shapes.add_textbox(c1_left + card_w - 550000, top_row_y + top_row_h - 550000, 350000, 350000)
link_1.text_frame.margin_left = 0
link_1.text_frame.margin_top = 0
p_l1 = link_1.text_frame.paragraphs[0]
p_l1.text = "🔗"
p_l1.font.size = Pt(11)

# CARD 2: Who faces it? (Top Middle)
c2_left = c1_left + card_w + col_gap
box_2 = slide_2.shapes.add_shape(5, c2_left, top_row_y, card_w, top_row_h)
box_2.fill.solid()
box_2.fill.fore_color.rgb = RGBColor(195, 106, 79)
box_2.line.fill.background()
# Title Box
t_box_2 = slide_2.shapes.add_textbox(c2_left + 150000, top_row_y + 80000, card_w - 300000, 350000)
t_box_2.text_frame.word_wrap = True
t_box_2.text_frame.margin_left = 0
t_box_2.text_frame.margin_top = 0
p_t2 = t_box_2.text_frame.paragraphs[0]
p_t2.text = "Who faces it?"
p_t2.font.name = "Arial"
p_t2.font.size = Pt(12)
p_t2.font.bold = True
p_t2.font.color.rgb = RGBColor(25, 25, 25)
# Inner White Card
w_card_2 = slide_2.shapes.add_shape(5, c2_left + 150000, top_row_y + 480000, card_w - 300000, top_row_h - 630000)
w_card_2.fill.solid()
w_card_2.fill.fore_color.rgb = RGBColor(255, 255, 255)
w_card_2.line.color.rgb = RGBColor(195, 106, 79)
w_card_2.line.width = Pt(1)
tf_2 = w_card_2.text_frame
tf_2.word_wrap = True
tf_2.margin_left = 120000
tf_2.margin_right = 120000
tf_2.margin_top = 100000
tf_2.margin_bottom = 100000
p_b2 = tf_2.paragraphs[0]
segments_2 = [
    ("Working professionals", True),
    (" using Claude for high-stakes knowledge work daily.\n\nOur Research:\n•  ", False),
    ("71%", True),
    (" of surveyed users rate their tasks as ", False),
    ("high-stakes", True),
    (".\n•  ", False),
    ("85%", True),
    (" interact with AI tools ", False),
    ("daily", True),
    (".", False)
]
add_formatted_text(p_b2, segments_2)

# CARD 3: How do we know it's real? (Top Right)
c3_left = c2_left + card_w + col_gap
box_3 = slide_2.shapes.add_shape(5, c3_left, top_row_y, card_w, top_row_h)
box_3.fill.solid()
box_3.fill.fore_color.rgb = RGBColor(195, 106, 79)
box_3.line.fill.background()
# Title Box
t_box_3 = slide_2.shapes.add_textbox(c3_left + 150000, top_row_y + 80000, card_w - 300000, 350000)
t_box_3.text_frame.word_wrap = True
t_box_3.text_frame.margin_left = 0
t_box_3.text_frame.margin_top = 0
p_t3 = t_box_3.text_frame.paragraphs[0]
p_t3.text = "How do we know it's real?"
p_t3.font.name = "Arial"
p_t3.font.size = Pt(12)
p_t3.font.bold = True
p_t3.font.color.rgb = RGBColor(25, 25, 25)
# Inner White Card
w_card_3 = slide_2.shapes.add_shape(5, c3_left + 150000, top_row_y + 480000, card_w - 300000, top_row_h - 630000)
w_card_3.fill.solid()
w_card_3.fill.fore_color.rgb = RGBColor(255, 255, 255)
w_card_3.line.color.rgb = RGBColor(195, 106, 79)
w_card_3.line.width = Pt(1)
tf_3 = w_card_3.text_frame
tf_3.word_wrap = True
tf_3.margin_left = 120000
tf_3.margin_right = 120000
tf_3.margin_top = 100000
tf_3.margin_bottom = 100000
p_b3 = tf_3.paragraphs[0]
segments_3 = [
    ("Users are trapped between blind trust and exhausting skepticism:\n•  ", False),
    ("30% Over-trust", True),
    (" (Fluency Bias Trap)\n•  ", False),
    ("37% Skeptic Trap", True),
    (" (re-checking endlessly)\n\n", False),
    ("*Grounded in: 35 surveys + 6 interviews + 2 task observations*", True)
]
add_formatted_text(p_b3, segments_3)

# CARD 4: What value does solving it generate? (Bottom Left - Wide Card)
c4_left = 457200
c4_width = card_w + col_gap + card_w # spans 2 columns
box_4 = slide_2.shapes.add_shape(5, c4_left, bottom_row_y, c4_width, bottom_row_h)
box_4.fill.solid()
box_4.fill.fore_color.rgb = RGBColor(195, 106, 79)
box_4.line.fill.background()
# Title Box
t_box_4 = slide_2.shapes.add_textbox(c4_left + 150000, bottom_row_y + 80000, c4_width - 300000, 350000)
t_box_4.text_frame.word_wrap = True
t_box_4.text_frame.margin_left = 0
t_box_4.text_frame.margin_top = 0
p_t4 = t_box_4.text_frame.paragraphs[0]
p_t4.text = "What value does solving it generate?"
p_t4.font.name = "Arial"
p_t4.font.size = Pt(12)
p_t4.font.bold = True
p_t4.font.color.rgb = RGBColor(25, 25, 25)
# Nested White Card 1: For Users
c4_inner_w = (c4_width - 300000 - 200000) / 2
w_card_4_left = slide_2.shapes.add_shape(5, c4_left + 150000, bottom_row_y + 480000, c4_inner_w, bottom_row_h - 630000)
w_card_4_left.fill.solid()
w_card_4_left.fill.fore_color.rgb = RGBColor(255, 255, 255)
w_card_4_left.line.color.rgb = RGBColor(195, 106, 79)
w_card_4_left.line.width = Pt(1)
tf_4_left = w_card_4_left.text_frame
tf_4_left.word_wrap = True
tf_4_left.margin_left = 120000
tf_4_left.margin_right = 120000
tf_4_left.margin_top = 100000
tf_4_left.margin_bottom = 100000
p_b4_l = tf_4_left.paragraphs[0]
segments_4_left = [
    ("For Users\n\n", True),
    ("•  ", False),
    ("Calibrated confidence", True),
    (" in outputs, removing the need to endlessly recheck stable results.\n\n•  Significant ", False),
    ("speedups in auditing", True),
    (" and risk avoidance on high-stakes tasks.", False)
]
add_formatted_text(p_b4_l, segments_4_left)

# Nested White Card 2: For Business
w_card_4_right = slide_2.shapes.add_shape(5, c4_left + 150000 + c4_inner_w + 200000, bottom_row_y + 480000, c4_inner_w, bottom_row_h - 630000)
w_card_4_right.fill.solid()
w_card_4_right.fill.fore_color.rgb = RGBColor(255, 255, 255)
w_card_4_right.line.color.rgb = RGBColor(195, 106, 79)
w_card_4_right.line.width = Pt(1)
tf_4_right = w_card_4_right.text_frame
tf_4_right.word_wrap = True
tf_4_right.margin_left = 120000
tf_4_right.margin_right = 120000
tf_4_right.margin_top = 100000
tf_4_right.margin_bottom = 100000
p_b4_r = tf_4_right.paragraphs[0]
segments_4_right = [
    ("For Anthropic\n\n", True),
    ("•  Drives long-term retention and stickiness of Claude in high-stakes enterprise workflows.\n\n•  Drives Pro & Team upgrades by demonstrating clear, auditable logic that justifies subscription premiums.", False)
]
add_formatted_text(p_b4_r, segments_4_right)
# Small Link Icon inside Card 4 Right
link_4 = slide_2.shapes.add_textbox(c4_left + 150000 + c4_inner_w + 200000 + c4_inner_w - 550000, bottom_row_y + bottom_row_h - 550000, 350000, 350000)
link_4.text_frame.margin_left = 0
link_4.text_frame.margin_top = 0
p_l4 = link_4.text_frame.paragraphs[0]
p_l4.text = "🔗"
p_l4.font.size = Pt(11)

# CARD 5: Why should we solve this now? (Bottom Right)
c5_left = c3_left
box_5 = slide_2.shapes.add_shape(5, c5_left, bottom_row_y, card_w, bottom_row_h)
box_5.fill.solid()
box_5.fill.fore_color.rgb = RGBColor(195, 106, 79)
box_5.line.fill.background()
# Title Box
t_box_5 = slide_2.shapes.add_textbox(c5_left + 150000, bottom_row_y + 80000, card_w - 300000, 350000)
t_box_5.text_frame.word_wrap = True
t_box_5.text_frame.margin_left = 0
t_box_5.text_frame.margin_top = 0
p_t5 = t_box_5.text_frame.paragraphs[0]
p_t5.text = "Why should we solve this now?"
p_t5.font.name = "Arial"
p_t5.font.size = Pt(12)
p_t5.font.bold = True
p_t5.font.color.rgb = RGBColor(25, 25, 25)
# Inner White Card
w_card_5 = slide_2.shapes.add_shape(5, c5_left + 150000, bottom_row_y + 480000, card_w - 300000, bottom_row_h - 630000)
w_card_5.fill.solid()
w_card_5.fill.fore_color.rgb = RGBColor(255, 255, 255)
w_card_5.line.color.rgb = RGBColor(195, 106, 79)
w_card_5.line.width = Pt(1)
tf_5 = w_card_5.text_frame
tf_5.word_wrap = True
tf_5.margin_left = 120000
tf_5.margin_right = 120000
tf_5.margin_top = 100000
tf_5.margin_bottom = 100000
p_b5 = tf_5.paragraphs[0]
segments_5 = [
    ("•  ", False),
    ("Trust is a habit", True),
    (". Every week users spend defaulting to blind copy-pasting or basic prompting reinforces bad habits.\n\n•  The longer we delay introducing a structured evaluation layer, the harder it becomes to shift user behaviors.", False)
]
add_formatted_text(p_b5, segments_5)
# Small Link Icon inside Card 5
link_5 = slide_2.shapes.add_textbox(c5_left + card_w - 550000, bottom_row_y + bottom_row_h - 550000, 350000, 350000)
link_5.text_frame.margin_left = 0
link_5.text_frame.margin_top = 0
p_l5 = link_5.text_frame.paragraphs[0]
p_l5.text = "🔗"
p_l5.font.size = Pt(11)

# ==========================================
# STEP 3: MODIFY SLIDE 6 (IDEATION / RICE TABLE)
# ==========================================
print("Modifying Slide 6 (index 5) RICE scores and text frames...")
slide_6 = prs.slides[5]

# 1. Update Idea Wordings (What it is / Promise / Limit)
slide_6.shapes[14].text_frame.text = "Claude infers the task type from the user's prompt and tailors both the answer and the review criteria to it (e.g. coding, research, or career)."
slide_6.shapes[16].text_frame.text = "Zero upfront friction. Keeps Claude's speed-preserving core while offering a post-answer review chip."
slide_6.shapes[18].text_frame.text = "Still depends on Claude correctly inferring the task lens, and the review format can feel slightly abstract."

slide_6.shapes[25].text_frame.text = "Surfaced only on ambiguous/high-stakes prompts. Claude previews 2–3 possible ways it could approach the task (e.g. fast, risk-aware)."
slide_6.shapes[27].text_frame.text = "Allows users to steer the response strategy before Claude commits, making the reasoning strategy visible."
slide_6.shapes[29].text_frame.text = "Still introduces a moment of hesitation before users see value. Risks feeling fragile in conversational speed."

slide_6.shapes[36].text_frame.text = "Claude answers first at full speed, then offers an optional evaluation layer where users route the answer through 6 specialized peer-evaluators (Code, Research, Writing, Reasoning, Career, Risk) and 3 playbooks (Balanced, Rigor, Style) powered by Llama 3.1 70B NIM models."
slide_6.shapes[38].text_frame.text = "Opt-in = no speed friction by default. Outputs transparent checklists showing strengths, weaknesses, and a concrete 'Verify before use' manual checklist. No opaque trust scores."
slide_6.shapes[40].text_frame.text = "Slightly higher engineering effort than simple review chips because it requires evaluator orchestration, parallel NIM routing, and UI safeguards."

text_box_indices = [14, 16, 18, 25, 27, 29, 36, 38, 40]
for idx in text_box_indices:
    tf = slide_6.shapes[idx].text_frame
    tf.word_wrap = True
    for p in tf.paragraphs:
        for run in p.runs:
            run.font.name = "Arial"
            run.font.size = Pt(10)
            run.font.color.rgb = RGBColor(240, 240, 240)

# 2. Update RICE Table Cell Values
rice_updates = {
    # Task Lens + Review Mode (Shape 50)
    51: "7.0",   # Reach
    52: "2.5",   # Impact
    53: "75%",   # Confidence
    54: "4.0",   # Effort
    55: "3.28",  # RICE Score
    
    # Interactive Decision Paths (Shape 57)
    58: "6.5",   # Reach
    59: "2.25",  # Impact
    60: "70%",   # Confidence
    61: "3.5",   # Effort
    62: "2.93",  # RICE Score
    
    # Progressive Eval + Evaluators (Shape 64)
    65: "8.5",   # Reach
    66: "3.0",   # Impact
    67: "85%",   # Confidence
    68: "4.5",   # Effort
    69: "4.82"   # RICE Score
}

print("Updating RICE Table cell shapes...")
for idx, new_val in rice_updates.items():
    shape = slide_6.shapes[idx]
    shape.text_frame.text = new_val
    for p in shape.text_frame.paragraphs:
        p.alignment = PP_ALIGN.CENTER
        for run in p.runs:
            run.font.name = "Arial"
            run.font.size = Pt(11)
            if idx in [65, 66, 67, 68, 69]:
                run.font.bold = True
                run.font.color.rgb = RGBColor(0, 180, 216) # Teal
            else:
                run.font.color.rgb = RGBColor(255, 255, 255) # White

# Save presentation
print("Saving modified presentation...")
prs.save("NL_Claude_TrustEvaluators.pptx")
print("PowerPoint file successfully updated in place and saved!")
