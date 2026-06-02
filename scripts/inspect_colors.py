import sys
from pptx import Presentation
from pptx.enum.dml import MSO_COLOR_TYPE

sys.stdout.reconfigure(encoding='utf-8')
prs = Presentation("NL_Claude_TrustEvaluators.pptx")

def print_shape_color(slide_idx, name, shape):
    fill = shape.fill
    print(f"Slide {slide_idx+1} - Shape '{shape.name}': Type={shape.shape_type}")
    if fill.type == 1: # SOLID
        print(f"  Solid Color: RGB={fill.fore_color.rgb if fill.fore_color.type == 1 else 'Not RGB'}")
    else:
        print(f"  Fill Type: {fill.type}")

print("--- Slide 1 ---")
for i, s in enumerate(prs.slides[0].shapes):
    if s.fill.type is not None:
        print_shape_color(0, s.name, s)

print("--- Slide 3 ---")
for i, s in enumerate(prs.slides[2].shapes):
    if s.fill.type is not None:
        print_shape_color(2, s.name, s)
